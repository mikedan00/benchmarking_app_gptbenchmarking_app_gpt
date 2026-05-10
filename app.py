#!/usr/bin/env python3
"""
AI 벤치마킹 보고서 자동 작성 시스템
HuggingFace LLM (google/gemma-4-26B-A4B-it) 기반 경쟁사 분석 통합 플랫폼
"""

import streamlit as st
import os
import json
import re
import time
import requests
import pandas as pd
from pathlib import Path
from io import BytesIO
from datetime import datetime
from urllib.parse import urlparse, urlunparse, parse_qsl, urlencode, quote_plus
from collections import Counter
import html
import PyPDF2
from docx import Document as DocxDocument
from bs4 import BeautifulSoup
from huggingface_hub import InferenceClient
from duckduckgo_search import DDGS
from dotenv import load_dotenv

load_dotenv()


def get_secret(name: str, default: str = "") -> str:
    """로컬 secrets.toml이 없어도 안전하게 st.secrets 값을 읽기 위한 헬퍼."""
    try:
        return st.secrets.get(name, default)
    except Exception:
        return default

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PAGE CONFIGURATION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
st.set_page_config(
    page_title="AI 벤치마킹 보고서 시스템",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
    menu_items={
        'About': "AI 기반 벤치마킹 보고서 자동 생성 시스템 | Powered by HuggingFace LLM"
    }
)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# CUSTOM CSS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
st.markdown("""
<style>
    .stApp { background: #0f172a; }
    .main .block-container { padding-top: 1rem; max-width: 1400px; }

    .hero-header {
        background: linear-gradient(135deg, #1e40af 0%, #7c3aed 50%, #0891b2 100%);
        border-radius: 16px; padding: 36px 32px; margin-bottom: 24px;
        text-align: center; box-shadow: 0 20px 60px rgba(124,58,237,0.4);
    }
    .hero-title { font-size: 2.4rem; font-weight: 800; color: #ffffff; margin: 0 0 8px 0; }
    .hero-subtitle { font-size: 1.05rem; color: rgba(255,255,255,0.85); margin: 0; }
    .hero-badge {
        display: inline-block; background: rgba(255,255,255,0.15);
        border: 1px solid rgba(255,255,255,0.3); border-radius: 999px;
        padding: 4px 14px; font-size: 0.78rem; color: rgba(255,255,255,0.9); margin-top: 14px;
    }
    .step-card {
        background: linear-gradient(135deg, #1e293b, #0f172a);
        border: 1px solid #334155; border-radius: 12px; padding: 20px; margin-bottom: 16px;
    }
    .step-card h4 { color: #e2e8f0; margin: 0 0 8px 0; font-size: 1rem; }
    .step-card p  { color: #94a3b8; margin: 0; font-size: 0.88rem; }

    .badge-ok   { background:#064e3b;color:#6ee7b7;border:1px solid #065f46;border-radius:999px;padding:3px 10px;font-size:0.75rem; }
    .badge-warn { background:#422006;color:#fcd34d;border:1px solid #92400e;border-radius:999px;padding:3px 10px;font-size:0.75rem; }

    .voc-card { background:#1e293b;border:1px solid #334155;border-left:4px solid #7c3aed;border-radius:8px;padding:14px;margin-bottom:10px; }
    .voc-card .voc-title { color:#c4b5fd;font-weight:600;font-size:0.9rem;margin-bottom:6px; }
    .voc-card .voc-body  { color:#94a3b8;font-size:0.82rem;line-height:1.5; }

    section[data-testid="stSidebar"] { background: linear-gradient(180deg,#0f172a 0%,#1e1b4b 100%) !important; border-right:1px solid #334155; }

    .stTabs [data-baseweb="tab-list"] { background:#1e293b;border-radius:10px;padding:4px;gap:4px; }
    .stTabs [data-baseweb="tab"]      { border-radius:8px;color:#64748b;font-weight:500; }
    .stTabs [aria-selected="true"]    { background:linear-gradient(135deg,#1e40af,#7c3aed) !important;color:white !important; }

    .stButton>button { background:linear-gradient(135deg,#1e40af,#7c3aed);color:white;border:none;border-radius:8px;font-weight:600; }
    .stProgress>div>div { background:linear-gradient(90deg,#1e40af,#7c3aed);border-radius:999px; }
</style>
""", unsafe_allow_html=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SESSION STATE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
DEFAULTS = {
    # st.secrets (Streamlit Cloud) → os.getenv (.env 로컬) → 빈 문자열 순으로 탐색
    "hf_token": (
        get_secret("HF_TOKEN", "")
        or os.getenv("HF_TOKEN", "")
    ),
    "model_id": (
        get_secret("HF_MODEL_ID", "")
        or os.getenv("HF_MODEL_ID", "")
        or "google/gemma-4-26B-A4B-it"
    ),
    "hf_provider": (
        get_secret("HF_PROVIDER", "")
        or os.getenv("HF_PROVIDER", "")
        or "auto"
    ),
    "benchmark_guide": None,
    "guide_extracted": False,
    "voc_my": {}, "voc_comp": {},
    "voc_analyzed_my": {}, "voc_analyzed_comp": {},
    "spec_comparison": {},
    "spec_my_raw": "", "spec_comp_raw": "",
    "report_md": "",
    "my_company": "", "my_product": "",
    "competitor": "", "comp_product": "",
}
for k, v in DEFAULTS.items():
    if k not in st.session_state:
        st.session_state[k] = v


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LLM CLIENT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def get_client():
    """HuggingFace Inference Providers 클라이언트 생성.

    중요:
    - Gemma/ZAYA 같은 instruct 모델은 대부분 conversational/chat-completion 경로로 호출해야 합니다.
    - provider="auto"는 HuggingFace가 사용 가능한 provider를 자동 선택합니다.
    """
    token = st.session_state.hf_token.strip() if st.session_state.hf_token else ""
    if not token:
        return None

    provider = (st.session_state.get("hf_provider") or "auto").strip()
    kwargs = {"token": token}
    if provider:
        kwargs["provider"] = provider

    try:
        return InferenceClient(**kwargs)
    except TypeError:
        # 구버전 huggingface_hub 호환: provider 인자를 모르는 경우
        return InferenceClient(token=token)


def _extract_chat_content(resp) -> str:
    """huggingface_hub 버전별 응답 객체/딕셔너리 차이를 흡수."""
    try:
        choice = resp.choices[0]
        msg = choice.message
        if isinstance(msg, dict):
            return (msg.get("content") or "").strip()
        return (getattr(msg, "content", "") or "").strip()
    except Exception:
        if isinstance(resp, dict):
            try:
                return (resp["choices"][0]["message"]["content"] or "").strip()
            except Exception:
                pass
    return str(resp).strip()


def call_llm(prompt: str, max_tokens: int = 2048, system: str = "") -> str:
    client = get_client()
    if client is None:
        return "⚠️ HuggingFace API 토큰을 먼저 사이드바에 입력해주세요."

    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})

    model_id = st.session_state.model_id.strip()
    provider = st.session_state.get("hf_provider", "auto")

    try:
        # 1순위: Inference Providers의 chat_completion API
        # Gemma 4 / ZAYA1 등 conversational 모델은 text_generation이 아니라 이 경로를 사용해야 합니다.
        if hasattr(client, "chat_completion"):
            resp = client.chat_completion(
                messages=messages,
                model=model_id,
                max_tokens=max_tokens,
                temperature=0.25,
            )
        else:
            # 일부 버전 호환: OpenAI-compatible chat.completions 경로
            resp = client.chat.completions.create(
                model=model_id,
                messages=messages,
                max_tokens=max_tokens,
                temperature=0.25,
            )
        content = _extract_chat_content(resp)
        return content if content else "⚠️ LLM 응답이 비어 있습니다. 모델/provider 조합을 변경해 보세요."
    except Exception as e:
        return (
            f"❌ LLM API 오류: {e}\n\n"
            "해결 방법:\n"
            "1) 이 앱은 conversational/chat-completion 방식으로 호출하도록 수정되어야 합니다.\n"
            "2) requirements.txt의 huggingface-hub를 최신 버전으로 업그레이드하세요.\n"
            "3) 사이드바에서 Provider를 auto 또는 해당 모델을 지원하는 provider로 바꿔 테스트하세요.\n"
            f"현재 설정: model={model_id}, provider={provider}"
        )


def safe_json(text: str):
    try:
        cleaned = re.sub(r"```(?:json)?", "", text).strip().rstrip("`").strip()
        start = min(
            (cleaned.find("{") if "{" in cleaned else len(cleaned)),
            (cleaned.find("[") if "[" in cleaned else len(cleaned)),
        )
        if start < len(cleaned):
            return json.loads(cleaned[start:])
    except Exception:
        pass
    return None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# FILE TEXT EXTRACTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def extract_pdf(data: bytes) -> str:
    try:
        reader = PyPDF2.PdfReader(BytesIO(data))
        return "\n".join(p.extract_text() or "" for p in reader.pages)
    except Exception as e:
        return f"[PDF 추출 오류: {e}]"

def extract_docx(data: bytes) -> str:
    try:
        doc = DocxDocument(BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[DOCX 추출 오류: {e}]"

def extract_file(uploaded) -> str:
    raw = uploaded.read()
    name = uploaded.name.lower()
    if name.endswith(".pdf"):
        return extract_pdf(raw)
    elif name.endswith((".docx", ".doc")):
        return extract_docx(raw)
    else:
        return raw.decode("utf-8", errors="ignore")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# DEFAULT GUIDE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
DEFAULT_GUIDE = {
    "title": "벤치마킹 보고서 표준 작성 가이드 v1.0",
    "summary": "경쟁사 비교 분석을 위한 체계적 벤치마킹 보고서 작성 기준",
    "main_sections": [
        {"id":"1","name":"경영진 요약 (Executive Summary)",
         "items":["핵심 발견사항 Top 5","전략적 시사점","즉각적 액션 아이템"],
         "criteria":["명확성","실행가능성","비즈니스 임팩트"]},
        {"id":"2","name":"분석 범위 및 방법론",
         "items":["분석 목적","대상 기업/제품","데이터 수집 방법","분석 기간"],
         "criteria":["객관성","데이터 신뢰성"]},
        {"id":"3","name":"기업 및 제품 개요 비교",
         "items":["회사 규모/역사","제품 포트폴리오","시장 포지셔닝","가격 전략"],
         "criteria":["시장점유율","성장률","브랜드 인지도"]},
        {"id":"4","name":"고객 VOC 분석",
         "items":["고객 만족도 CSAT","NPS 비교","주요 불만 사항","칭찬 포인트","개선 요청"],
         "criteria":["감성 분석","빈도 분석","영향도 분석"]},
        {"id":"5","name":"제품/서비스 사양 및 성능 비교",
         "items":["핵심 기능 매트릭스","성능 벤치마크","기술 사양","가격 대비 성능"],
         "criteria":["기능 완성도","성능 우위","혁신성"]},
        {"id":"6","name":"차별화 강점·약점·갭 분석",
         "items":["고유 기능 분석","SWOT 매트릭스","기능 갭","기회 영역"],
         "criteria":["차별화 지속 가능성","모방 난이도"]},
        {"id":"7","name":"전략적 권고사항",
         "items":["단기 액션 0~6M","중기 전략 6~18M","장기 로드맵 18M+"],
         "criteria":["ROI 예상","실행 난이도","우선순위"]},
    ],
    "scoring_dimensions": ["기능성","성능","가격경쟁력","고객만족도","혁신성","지원/서비스"],
    "scoring_scale": "1~10점 척도 (10점: 압도적 우위)",
    "data_sources": ["웹 리뷰 (블로그/커뮤니티)","공식 제품 스펙 페이지","유통 채널 리뷰","SNS VOC","IT 전문 미디어"],
    "analysis_methods": ["감성 분석","갭 분석","SWOT 분석","기능 매트릭스","가치 맵"],
}


def build_guide_from_files(texts: list) -> dict:
    combined = "\n\n━━ 문서 구분 ━━\n\n".join(texts[:15])
    if len(combined) > 10_000:
        combined = combined[:10_000] + "\n...[토큰 제한으로 생략됨]"

    system = "당신은 벤치마킹 보고서 전문 컨설턴트입니다. 여러 보고서를 분석하여 표준 가이드를 JSON으로 추출합니다. 한국어로 답하고, 반드시 순수 JSON만 출력하세요."
    prompt = f"""아래는 여러 벤치마킹 보고서의 내용입니다. 공통 패턴을 분석해 표준 작성 가이드를 추출하세요.

[보고서 내용]
{combined}

다음 JSON 스키마로 정확히 출력하세요 (주석·마크다운 없이 JSON만):
{{
  "title": "벤치마킹 보고서 표준 작성 가이드 v1.0",
  "summary": "한 문장 요약",
  "main_sections": [
    {{"id":"1","name":"섹션명","items":["핵심항목1","핵심항목2"],"criteria":["평가기준1"]}}
  ],
  "scoring_dimensions": ["차원1","차원2"],
  "scoring_scale": "평가 척도 설명",
  "data_sources": ["출처1","출처2"],
  "analysis_methods": ["방법1","방법2"]
}}"""

    raw = call_llm(prompt, max_tokens=3000, system=system)
    parsed = safe_json(raw)
    if parsed and isinstance(parsed, dict) and "main_sections" in parsed:
        return parsed
    guide = DEFAULT_GUIDE.copy()
    guide["llm_notes"] = raw
    return guide


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# WEB SEARCH & VOC
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
VOC_CATEGORY_KEYWORDS = {
    "성능/속도": ["성능", "속도", "버벅", "렉", "느림", "딜레이", "최적화", "performance", "lag", "slow"],
    "배터리/발열": ["배터리", "방전", "충전", "발열", "뜨거", "온도", "battery", "charging", "heat", "overheat"],
    "카메라/화질": ["카메라", "사진", "동영상", "줌", "화질", "야간", "camera", "photo", "video", "zoom"],
    "디스플레이/디자인": ["디스플레이", "화면", "밝기", "주사율", "무게", "디자인", "display", "screen", "design", "weight"],
    "가격/가성비": ["가격", "비싸", "가성비", "할인", "요금", "price", "expensive", "value", "cost"],
    "소프트웨어/UI": ["UI", "UX", "업데이트", "버그", "앱", "소프트웨어", "software", "update", "bug", "crash"],
    "내구성/품질": ["고장", "불량", "파손", "스크래치", "품질", "내구", "defect", "durability", "broken"],
    "A/S/고객지원": ["AS", "A/S", "서비스센터", "교환", "환불", "상담", "support", "warranty", "refund"],
    "호환성/생태계": ["호환", "연동", "액세서리", "생태계", "compatibility", "ecosystem", "accessory"],
    "보안/개인정보": ["보안", "개인정보", "프라이버시", "security", "privacy"],
}

POSITIVE_WORDS = [
    "좋", "만족", "추천", "훌륭", "빠르", "선명", "편하", "예쁘", "강력", "개선", "최고",
    "good", "great", "excellent", "satisfied", "recommend", "fast", "clear", "better", "best"
]
NEGATIVE_WORDS = [
    "불만", "문제", "별로", "나쁘", "실망", "느리", "버벅", "고장", "불량", "비싸", "아쉽", "최악", "발열",
    "bad", "issue", "problem", "poor", "slow", "defect", "broken", "expensive", "disappoint", "overheat", "bug"
]
REQUEST_WORDS = ["개선", "요청", "바라", "필요", "해결", "추가", "지원", "원함", "should", "need", "wish", "request", "fix", "improve"]


def clean_text(text: str) -> str:
    text = html.unescape(str(text or ""))
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def normalize_url(url: str) -> str:
    """URL 중복 제거용 정규화: UTM/fbclid/gclid 등 추적 파라미터 제거."""
    if not url:
        return ""
    try:
        parsed = urlparse(url)
        query = [
            (k, v) for k, v in parse_qsl(parsed.query, keep_blank_values=True)
            if not (k.lower().startswith("utm_") or k.lower() in {"fbclid", "gclid", "mc_cid", "mc_eid", "igshid"})
        ]
        netloc = parsed.netloc.lower().replace("www.", "")
        path = parsed.path.rstrip("/")
        return urlunparse((parsed.scheme or "https", netloc, path, "", urlencode(query), ""))
    except Exception:
        return url.strip()


def get_domain(url: str) -> str:
    try:
        return urlparse(url).netloc.lower().replace("www.", "")
    except Exception:
        return ""


def infer_channel(domain: str, title: str = "") -> str:
    blob = f"{domain} {title}".lower()
    if any(x in blob for x in ["youtube", "youtu.be"]):
        return "동영상"
    if any(x in blob for x in ["news", "zdnet", "etnews", "chosun", "hankyung", "mk.co.kr", "yna", "theverge", "engadget"]):
        return "뉴스/미디어"
    if any(x in blob for x in ["blog", "tistory", "naver.com", "medium"]):
        return "블로그/후기"
    if any(x in blob for x in ["reddit", "quora", "dcinside", "clien", "ppomppu", "ruliweb", "fmkorea", "community", "forum"]):
        return "커뮤니티"
    if any(x in blob for x in ["coupang", "amazon", "bestbuy", "danawa", "11st", "gmarket", "shopping"]):
        return "쇼핑/리뷰"
    return "웹"


def infer_voc_category(text: str) -> str:
    t = text.lower()
    scores = {}
    for cat, words in VOC_CATEGORY_KEYWORDS.items():
        scores[cat] = sum(t.count(w.lower()) for w in words)
    best_cat, best_score = max(scores.items(), key=lambda x: x[1])
    return best_cat if best_score > 0 else "기타"


def infer_sentiment(text: str) -> tuple[str, int]:
    t = text.lower()
    pos = sum(t.count(w.lower()) for w in POSITIVE_WORDS)
    neg = sum(t.count(w.lower()) for w in NEGATIVE_WORDS)
    req = sum(t.count(w.lower()) for w in REQUEST_WORDS)
    score = pos - neg
    if neg >= 2 and pos >= 2:
        label = "혼합"
    elif neg > pos:
        label = "부정"
    elif pos > neg:
        label = "긍정"
    elif req > 0:
        label = "개선요청"
    else:
        label = "중립"
    return label, score


def build_voc_queries(company: str, product: str, depth: str = "standard") -> list[str]:
    target = clean_text(f"{company} {product}").strip()
    base_queries = [
        f"{target} 사용후기 리뷰 장단점",
        f"{target} VOC 불만 개선사항",
        f"{target} 문제점 불량 이슈",
        f'"{company}" "{product}" review pros cons',
    ]
    if depth in {"standard", "deep"}:
        base_queries += [
            f"{target} 배터리 발열 카메라 성능 후기",
            f"{target} 업데이트 버그 AS 서비스센터",
            f"{target} 가격 가성비 만족도",
            f"{target} 커뮤니티 후기 클리앙 뽐뿌 reddit",
        ]
    if depth == "deep":
        base_queries += [
            f"{target} 단점 아쉬운점 사지마",
            f"{target} defect issue complaint",
            f"{target} battery camera display software problem",
            f"{target} customer review forum reddit",
        ]
    # 빈 제품명일 때 중복 쿼리 제거
    return list(dict.fromkeys(q for q in base_queries if q.strip()))


def ddg_search(query: str, n: int = 8, region: str = "kr-kr", timelimit: str | None = None) -> list:
    try:
        with DDGS() as ddgs:
            kwargs = {"max_results": n, "region": region, "safesearch": "moderate"}
            if timelimit:
                kwargs["timelimit"] = timelimit
            results = list(ddgs.text(query, **kwargs))
        time.sleep(0.4)
        return results
    except TypeError:
        try:
            with DDGS() as ddgs:
                results = list(ddgs.text(query, max_results=n, region=region))
            time.sleep(0.4)
            return results
        except Exception as e:
            st.warning(f"검색 실패 ({query[:40]}…): {e}")
            return []
    except Exception as e:
        st.warning(f"검색 실패 ({query[:40]}…): {e}")
        return []


def ddg_news_search(query: str, n: int = 8, region: str = "kr-kr", timelimit: str | None = "m") -> list:
    try:
        with DDGS() as ddgs:
            results = list(ddgs.news(query, max_results=n, region=region, safesearch="moderate", timelimit=timelimit))
        time.sleep(0.4)
        return [
            {
                "title": r.get("title", ""),
                "href": r.get("url") or r.get("href", ""),
                "body": r.get("body") or r.get("excerpt", ""),
                "date": r.get("date", ""),
                "source_type": "DDG News",
            }
            for r in results
        ]
    except Exception as e:
        st.warning(f"뉴스 검색 실패 ({query[:40]}…): {e}")
        return []


def google_news_rss(query: str, n: int = 8) -> list:
    """별도 API 키 없이 Google News RSS 검색 결과 수집."""
    url = f"https://news.google.com/rss/search?q={quote_plus(query)}&hl=ko&gl=KR&ceid=KR:ko"
    headers = {"User-Agent": "Mozilla/5.0"}
    try:
        res = requests.get(url, headers=headers, timeout=8)
        res.raise_for_status()
        soup = BeautifulSoup(res.text, "xml")
        items = []
        for it in soup.find_all("item")[:n]:
            items.append({
                "title": clean_text(it.title.text if it.title else ""),
                "href": it.link.text if it.link else "",
                "body": clean_text(it.description.text if it.description else ""),
                "date": clean_text(it.pubDate.text if it.pubDate else ""),
                "source_type": "Google News RSS",
            })
        time.sleep(0.3)
        return items
    except Exception as e:
        st.warning(f"Google News RSS 실패 ({query[:40]}…): {e}")
        return []


def fetch_page_excerpt(url: str, max_chars: int = 1600) -> str:
    """검색 snippet이 너무 짧을 때 본문 일부만 가볍게 보강."""
    if not url or re.search(r"\.(pdf|zip|docx?|xlsx?|pptx?)(\?|$)", url, re.I):
        return ""
    headers = {"User-Agent": "Mozilla/5.0"}
    try:
        res = requests.get(url, headers=headers, timeout=6)
        ctype = res.headers.get("content-type", "")
        if "text/html" not in ctype and "application/xhtml" not in ctype:
            return ""
        soup = BeautifulSoup(res.text[:350_000], "lxml")
        for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
            tag.decompose()
        article = soup.find("article") or soup.find("main") or soup.body
        paragraphs = [clean_text(p.get_text(" ")) for p in article.find_all("p")[:12]] if article else []
        text = " ".join(p for p in paragraphs if len(p) > 30)
        return text[:max_chars]
    except Exception:
        return ""


def enrich_voc_items(raw_items: list, company: str, product: str, fetch_pages: bool = False) -> list:
    seen_url, seen_title, enriched = set(), set(), []
    target_terms = [x for x in [company, product] if x]
    for item in raw_items:
        href = item.get("href") or item.get("url") or ""
        norm = normalize_url(href)
        title = clean_text(item.get("title", ""))
        body = clean_text(item.get("body") or item.get("snippet") or item.get("excerpt") or "")
        title_key = re.sub(r"\W+", "", title.lower())[:80]
        if not title and not body:
            continue
        if norm and norm in seen_url:
            continue
        if title_key and title_key in seen_title:
            continue
        if norm:
            seen_url.add(norm)
        if title_key:
            seen_title.add(title_key)

        excerpt = fetch_page_excerpt(href) if fetch_pages else ""
        full_text = clean_text(" ".join([title, body, excerpt]))
        domain = get_domain(href)
        category = infer_voc_category(full_text)
        sentiment, polarity = infer_sentiment(full_text)
        relevance = sum(full_text.lower().count(t.lower()) for t in target_terms if t)
        relevance += 2 if category != "기타" else 0
        relevance += 1 if sentiment in {"부정", "개선요청", "혼합"} else 0
        relevance += min(len(full_text) // 300, 3)

        enriched.append({
            "title": title,
            "href": href,
            "normalized_url": norm,
            "source_domain": domain,
            "channel": infer_channel(domain, title),
            "source_type": item.get("source_type", "DDG Web"),
            "date": item.get("date", ""),
            "body": body,
            "excerpt": excerpt,
            "text": full_text[:2400],
            "category": category,
            "sentiment_hint": sentiment,
            "polarity_score": polarity,
            "relevance_score": relevance,
            "collected_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        })
    enriched.sort(key=lambda x: (x["relevance_score"], len(x.get("text", ""))), reverse=True)
    return enriched


def collect_voc(
    company: str,
    product: str,
    n_per_query: int = 8,
    sources: list[str] | None = None,
    depth: str = "standard",
    fetch_pages: bool = False,
    region: str = "kr-kr",
    timelimit: str | None = None,
) -> list:
    """멀티소스 VOC 자동 수집.

    sources: ["web", "news", "rss"] 중 선택
    depth: basic / standard / deep
    timelimit: d=1일, w=1주, m=1개월, y=1년. None은 전체 기간.
    """
    sources = sources or ["web", "news", "rss"]
    queries = build_voc_queries(company, product, depth)
    raw = []
    for q in queries:
        if "web" in sources:
            raw.extend({**r, "source_type": "DDG Web", "query": q} for r in ddg_search(q, n_per_query, region=region, timelimit=timelimit))
        if "news" in sources:
            raw.extend({**r, "query": q} for r in ddg_news_search(q, max(3, n_per_query // 2), region=region, timelimit=timelimit or "m"))
        if "rss" in sources:
            raw.extend({**r, "query": q} for r in google_news_rss(q, max(3, n_per_query // 2)))
    return enrich_voc_items(raw, company, product, fetch_pages=fetch_pages)


def summarize_voc_locally(items: list) -> dict:
    category_counter = Counter(x.get("category", "기타") for x in items)
    sentiment_counter = Counter(x.get("sentiment_hint", "중립") for x in items)
    channel_counter = Counter(x.get("channel", "웹") for x in items)
    negative_items = [x for x in items if x.get("sentiment_hint") in {"부정", "혼합", "개선요청"}]
    positive_items = [x for x in items if x.get("sentiment_hint") == "긍정"]
    return {
        "voc_count": len(items),
        "category_distribution": dict(category_counter.most_common()),
        "sentiment_distribution": dict(sentiment_counter.most_common()),
        "channel_distribution": dict(channel_counter.most_common()),
        "top_issue_categories": [k for k, _ in category_counter.most_common(5)],
        "sample_negative_titles": [x.get("title", "") for x in negative_items[:5]],
        "sample_positive_titles": [x.get("title", "") for x in positive_items[:5]],
    }


def collect_specs(company: str, product: str) -> list:
    queries = [
        f"{company} {product} 제품 사양 스펙 specification",
        f"{company} {product} 기능 비교 성능",
        f"{company} {product} official features datasheet",
    ]
    results = []
    for q in queries:
        results.extend(ddg_search(q, 6))
    return results


def analyze_voc(company: str, product: str, items: list) -> dict:
    if not items:
        return {"error": "VOC 데이터 없음", "company": company, "product": product, "voc_count": 0}

    local_summary = summarize_voc_locally(items)
    snippets = "\n".join(
        f"[{i+1}] 제목: {r.get('title','')}\n"
        f"채널: {r.get('channel','')} / 출처: {r.get('source_domain','')} / 분류: {r.get('category','')} / 감성힌트: {r.get('sentiment_hint','')}\n"
        f"내용: {(r.get('text') or r.get('body') or '')[:700]}\n"
        f"URL: {r.get('href','')}"
        for i, r in enumerate(items[:35])
    )

    system = "당신은 고객 경험(CX), 제품기획, VOC 데이터 분석 전문가입니다. 반드시 순수 JSON으로만 답변하세요."
    prompt = f"""다음은 '{company} {product}'에 대한 웹/뉴스/커뮤니티 기반 VOC 후보 데이터입니다.
검색 결과에는 홍보성 글과 뉴스가 섞일 수 있으므로, 실제 사용자 불만·칭찬·개선 요구를 우선 추출하세요.

[로컬 1차 통계]
{json.dumps(local_summary, ensure_ascii=False, indent=2)}

[VOC 후보 원문]
{snippets}

다음 JSON 스키마로 분석 결과를 출력하세요:
{{
  "company": "{company}",
  "product": "{product}",
  "voc_count": {len(items)},
  "overall_sentiment": "긍정/부정/보통/혼합",
  "satisfaction_score": 7.5,
  "confidence": "상/중/하",
  "category_distribution": {{"카테고리": 0}},
  "sentiment_distribution": {{"긍정": 0, "중립": 0, "부정": 0, "개선요청": 0}},
  "top_positives": ["강점1", "강점2", "강점3"],
  "top_negatives": ["약점1", "약점2", "약점3"],
  "key_features_praised": ["칭찬기능1", "기능2"],
  "key_features_criticized": ["비판기능1", "기능2"],
  "improvement_requests": ["개선요청1", "요청2"],
  "trending_topics": ["이슈1", "이슈2"],
  "issue_matrix": [
    {{"issue":"이슈명", "category":"분류", "severity":"상/중/하", "frequency":"상/중/하", "business_impact":"영향", "recommended_action":"권고 액션"}}
  ],
  "competitor_implication": "경쟁 벤치마킹 관점의 시사점",
  "notable_snippets": ["핵심 근거 요약1", "핵심 근거 요약2"]
}}"""

    raw = call_llm(prompt, max_tokens=2400, system=system)
    result = safe_json(raw)
    if isinstance(result, dict):
        # LLM이 분포값을 비우거나 이상하게 줄 때 로컬 통계로 보강
        result.setdefault("category_distribution", local_summary["category_distribution"])
        result.setdefault("sentiment_distribution", local_summary["sentiment_distribution"])
        result.setdefault("channel_distribution", local_summary["channel_distribution"])
        result.setdefault("voc_count", len(items))
        return result

    # LLM 실패 시에도 리포트가 비지 않도록 로컬 통계 기반 fallback 반환
    return {
        "company": company,
        "product": product,
        "voc_count": len(items),
        "overall_sentiment": max(local_summary["sentiment_distribution"], key=local_summary["sentiment_distribution"].get) if items else "N/A",
        "satisfaction_score": "N/A",
        "confidence": "하",
        "category_distribution": local_summary["category_distribution"],
        "sentiment_distribution": local_summary["sentiment_distribution"],
        "channel_distribution": local_summary["channel_distribution"],
        "top_positives": local_summary["sample_positive_titles"],
        "top_negatives": local_summary["sample_negative_titles"],
        "improvement_requests": [],
        "trending_topics": local_summary["top_issue_categories"],
        "raw": raw,
    }


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SPEC COMPARISON
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def compare_specs(my_co, my_prod, comp_co, comp_prod, my_text, comp_text) -> dict:
    system = "당신은 IT 제품 분석 전문가입니다. 두 제품을 객관적으로 비교해 순수 JSON으로만 출력하세요."
    prompt = f"""다음 두 제품의 정보를 비교 분석해주세요.

[자사: {my_co} – {my_prod}]
{my_text[:2500] if my_text else '직접 입력된 사양 없음 (웹 수집 데이터 활용)'}

[경쟁사: {comp_co} – {comp_prod}]
{comp_text[:2500] if comp_text else '직접 입력된 사양 없음 (웹 수집 데이터 활용)'}

다음 JSON으로 출력하세요:
{{
  "overall_summary": "한 문장 종합 비교",
  "overall_winner": "{my_co} 우위 / {comp_co} 우위 / 동등",
  "spec_matrix": [
    {{"category":"카테고리","my_value":"자사 사양","comp_value":"경쟁사 사양","winner":"{my_co}/{comp_co}/동등","importance":"상/중/하","note":"비고"}}
  ],
  "my_strengths": ["자사강점1","강점2","강점3"],
  "comp_strengths": ["경쟁사강점1","강점2","강점3"],
  "my_unique_only": ["자사고유기능1","기능2"],
  "comp_unique_only": ["경쟁사고유기능1","기능2"],
  "my_gaps": ["자사부족영역1","영역2"],
  "comp_gaps": ["경쟁사부족영역1","영역2"],
  "differentiation_score": {{
    "functionality":   {{"my":7,"comp":8}},
    "performance":     {{"my":8,"comp":7}},
    "price_value":     {{"my":7,"comp":6}},
    "innovation":      {{"my":8,"comp":9}},
    "customer_support":{{"my":7,"comp":7}}
  }},
  "strategic_recommendation": "전략적 권고사항"
}}"""

    raw = call_llm(prompt, max_tokens=3000, system=system)
    result = safe_json(raw)
    return result if isinstance(result, dict) else {"raw": raw}


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# REPORT GENERATION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def generate_report(guide, my_co, my_prod, comp_co, comp_prod, my_voc, comp_voc, spec_comp) -> str:
    guide_str  = json.dumps(guide,     ensure_ascii=False)[:1800]
    my_voc_s   = json.dumps(my_voc,    ensure_ascii=False)[:1000]
    comp_voc_s = json.dumps(comp_voc,  ensure_ascii=False)[:1000]
    spec_s     = json.dumps(spec_comp, ensure_ascii=False)[:2000]
    today      = datetime.now().strftime("%Y년 %m월 %d일")

    system = (
        "당신은 McKinsey 수준의 전략 컨설턴트입니다. "
        "데이터 기반의 통찰력 있는 벤치마킹 보고서를 한국어 마크다운으로 작성합니다. "
        "구체적 수치, 비교표, 실행 가능한 권고사항을 반드시 포함하세요."
    )
    prompt = f"""아래 데이터를 바탕으로 완전한 벤치마킹 보고서를 작성하세요.

## 메타 정보
- 자사: {my_co} ({my_prod})
- 경쟁사: {comp_co} ({comp_prod})
- 보고서 작성일: {today}

## 작성 가이드
{guide_str}

## VOC 분석 데이터
- 자사 VOC: {my_voc_s}
- 경쟁사 VOC: {comp_voc_s}

## 제품 사양 비교 데이터
{spec_s}

---

아래 구조로 전문 보고서를 작성하세요 (마크다운 형식):

# 📊 벤치마킹 보고서: {my_co} vs {comp_co}
**작성일:** {today} | **대상 제품:** {my_prod} vs {comp_prod}

---

## 1. 경영진 요약 (Executive Summary)
> 핵심 발견사항 5가지와 전략적 시사점

## 2. 분석 개요
## 3. 기업/제품 개요 비교 (마크다운 테이블)
## 4. 고객 VOC 심층 분석
### 4.1 {my_co} 고객 반응
### 4.2 {comp_co} 고객 반응
### 4.3 VOC 비교 시사점
## 5. 제품 사양 및 성능 상세 비교 (매트릭스 테이블)
## 6. 차별화 강점·약점·갭 분석
### 6.1 {my_co} 고유 강점
### 6.2 {my_co} 취약 영역 및 갭
### 6.3 {comp_co} 고유 강점
### 6.4 기능 갭 요약표
## 7. SWOT 분석 (2x2 마크다운 테이블)
## 8. 전략적 권고사항
### 8.1 단기 액션 플랜 (0~6개월)
### 8.2 중기 전략 (6~18개월)
### 8.3 장기 로드맵 (18개월+)
## 9. 결론

---
*본 보고서는 AI 벤치마킹 시스템으로 자동 생성되었습니다.*
*생성 모델: {st.session_state.model_id}*"""

    return call_llm(prompt, max_tokens=4096, system=system)




# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MS OFFICE EXPORTERS: DOCX / PPTX
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def safe_filename(value: str, fallback: str = "benchmarking_report") -> str:
    value = clean_text(value or "")
    value = re.sub(r"[^0-9A-Za-z가-힣._-]+", "_", value).strip("._-")
    return value[:120] or fallback


def strip_md_inline(text: str) -> str:
    """다운로드 문서용 간단 마크다운 인라인 제거."""
    text = html.unescape(str(text or ""))
    text = re.sub(r"!\[[^\]]*\]\([^\)]*\)", "", text)
    text = re.sub(r"\[([^\]]+)\]\([^\)]*\)", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"__([^_]+)__", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"_([^_]+)_", r"\1", text)
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.I)
    text = re.sub(r"<[^>]+>", "", text)
    return re.sub(r"\s+", " ", text).strip()


def parse_md_table(lines: list[str], start_idx: int):
    """마크다운 테이블을 (headers, rows, next_idx)로 파싱. 아니면 None 반환."""
    if start_idx + 1 >= len(lines):
        return None
    line = lines[start_idx].strip()
    sep = lines[start_idx + 1].strip()
    if "|" not in line or "|" not in sep:
        return None
    if not re.match(r"^\s*\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?\s*$", sep):
        return None

    def split_row(row: str) -> list[str]:
        row = row.strip().strip("|")
        return [strip_md_inline(c.strip()) for c in row.split("|")]

    headers = split_row(line)
    rows = []
    i = start_idx + 2
    while i < len(lines) and "|" in lines[i]:
        if lines[i].strip():
            row = split_row(lines[i])
            if len(row) < len(headers):
                row += [""] * (len(headers) - len(row))
            rows.append(row[:len(headers)])
        i += 1
    return headers, rows, i


def add_docx_korean_font(run, size_pt: int | None = None, bold: bool | None = None):
    try:
        from docx.oxml.ns import qn
        run.font.name = "Malgun Gothic"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Malgun Gothic")
        if size_pt:
            from docx.shared import Pt
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.bold = bold
    except Exception:
        pass


def build_docx_report_bytes(report_md: str, metadata: dict | None = None) -> bytes:
    """마크다운 보고서를 Word(.docx) 파일 bytes로 변환."""
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    metadata = metadata or {}
    doc = DocxDocument()

    # 기본 스타일
    try:
        normal = doc.styles["Normal"]
        normal.font.name = "Malgun Gothic"
        normal.font.size = Pt(10.5)
        for sty_name, size in [("Title", 22), ("Heading 1", 18), ("Heading 2", 15), ("Heading 3", 13)]:
            sty = doc.styles[sty_name]
            sty.font.name = "Malgun Gothic"
            sty.font.size = Pt(size)
            sty.font.color.rgb = RGBColor(30, 64, 175)
    except Exception:
        pass

    title = metadata.get("title") or "AI 벤치마킹 보고서"
    subtitle = metadata.get("subtitle") or ""
    meta_line = metadata.get("meta_line") or f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p_title.add_run(title)
    add_docx_korean_font(r, 22, True)

    if subtitle:
        p_sub = doc.add_paragraph()
        p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p_sub.add_run(subtitle)
        add_docx_korean_font(r, 12, False)

    p_meta = doc.add_paragraph()
    p_meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p_meta.add_run(meta_line)
    add_docx_korean_font(r, 9, False)
    doc.add_paragraph("")

    lines = (report_md or "").splitlines()
    i = 0
    while i < len(lines):
        raw = lines[i]
        line = raw.strip()

        if not line:
            i += 1
            continue

        parsed_table = parse_md_table(lines, i)
        if parsed_table:
            headers, rows, next_i = parsed_table
            table = doc.add_table(rows=1, cols=max(1, len(headers)))
            table.style = "Table Grid"
            for c, head in enumerate(headers):
                cell = table.rows[0].cells[c]
                cell.text = head
                for para in cell.paragraphs:
                    for run in para.runs:
                        add_docx_korean_font(run, 9, True)
            for row in rows[:80]:
                cells = table.add_row().cells
                for c, val in enumerate(row):
                    cells[c].text = val
                    for para in cells[c].paragraphs:
                        for run in para.runs:
                            add_docx_korean_font(run, 9, False)
            doc.add_paragraph("")
            i = next_i
            continue

        if re.match(r"^---+$", line):
            doc.add_paragraph("────────────────────────────────────────")
            i += 1
            continue

        m = re.match(r"^(#{1,4})\s+(.+)$", line)
        if m:
            level = min(len(m.group(1)), 3)
            txt = strip_md_inline(m.group(2))
            p = doc.add_heading(txt, level=level)
            for run in p.runs:
                add_docx_korean_font(run, 18 if level == 1 else 15 if level == 2 else 13, True)
            i += 1
            continue

        if line.startswith(">"):
            p = doc.add_paragraph(strip_md_inline(line.lstrip("> ")), style="Intense Quote")
            for run in p.runs:
                add_docx_korean_font(run, 10, False)
            i += 1
            continue

        if re.match(r"^[-*+]\s+", line):
            p = doc.add_paragraph(strip_md_inline(re.sub(r"^[-*+]\s+", "", line)), style="List Bullet")
            for run in p.runs:
                add_docx_korean_font(run, 10, False)
            i += 1
            continue

        if re.match(r"^\d+[.)]\s+", line):
            p = doc.add_paragraph(strip_md_inline(re.sub(r"^\d+[.)]\s+", "", line)), style="List Number")
            for run in p.runs:
                add_docx_korean_font(run, 10, False)
            i += 1
            continue

        p = doc.add_paragraph(strip_md_inline(line))
        for run in p.runs:
            add_docx_korean_font(run, 10, False)
        i += 1

    # Footer
    try:
        section = doc.sections[0]
        footer = section.footer.paragraphs[0]
        footer.text = "AI 벤치마킹 보고서 자동 작성 시스템"
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in footer.runs:
            add_docx_korean_font(run, 8, False)
    except Exception:
        pass

    bio = BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.getvalue()


def split_report_sections(report_md: str) -> list[tuple[str, list[str]]]:
    """PPT 슬라이드용으로 H2 기준 섹션 분리."""
    sections = []
    current_title = "핵심 요약"
    current_lines: list[str] = []
    for raw in (report_md or "").splitlines():
        line = raw.strip()
        if not line or re.match(r"^---+$", line):
            continue
        m = re.match(r"^##\s+(.+)$", line)
        if m:
            if current_lines:
                sections.append((strip_md_inline(current_title), current_lines))
            current_title = strip_md_inline(m.group(1))
            current_lines = []
            continue
        if line.startswith("# "):
            continue
        if "|" in line and re.search(r"\|\s*:?-{3,}:?", line):
            continue
        if "|" in line and len(line) > 20:
            # 표는 별도 데이터 또는 DOCX에서 처리하고 PPT에는 과밀 방지 목적으로 생략
            continue
        clean = strip_md_inline(re.sub(r"^(#{3,4})\s+", "", line))
        clean = re.sub(r"^[-*+]\s+", "", clean)
        clean = re.sub(r"^\d+[.)]\s+", "", clean)
        if clean and len(clean) > 2:
            current_lines.append(clean)
    if current_lines:
        sections.append((strip_md_inline(current_title), current_lines))
    return sections


def compact_bullets(lines: list[str], max_items: int = 7, max_chars: int = 115) -> list[str]:
    bullets = []
    seen = set()
    for line in lines:
        line = strip_md_inline(line)
        if not line or line in seen:
            continue
        seen.add(line)
        if len(line) > max_chars:
            line = line[:max_chars - 1].rstrip() + "…"
        bullets.append(line)
        if len(bullets) >= max_items:
            break
    return bullets


def add_pptx_bullet_slide(prs, title: str, bullets: list[str], accent_rgb=(30, 64, 175)):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 상단 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title[:64]
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.name = "Malgun Gothic"
    r.font.color.rgb = RGBColor(*accent_rgb)

    body_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.25), Inches(11.85), Inches(5.65))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets or ["요약 항목 없음"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18 if len(bullet) < 80 else 16)
        p.font.name = "Malgun Gothic"
        p.space_after = Pt(10)

    return slide


def build_pptx_report_bytes(report_md: str, metadata: dict | None = None, voc_my: dict | None = None, voc_comp: dict | None = None, spec_comp: dict | None = None) -> bytes:
    """마크다운 보고서를 PowerPoint(.pptx) 요약 덱으로 변환."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    metadata = metadata or {}
    voc_my = voc_my or {}
    voc_comp = voc_comp or {}
    spec_comp = spec_comp or {}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(15, 23, 42)
    title = metadata.get("title") or "AI 벤치마킹 보고서"
    subtitle = metadata.get("subtitle") or ""
    meta_line = metadata.get("meta_line") or datetime.now().strftime("%Y-%m-%d %H:%M")

    tb = slide.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(11.7), Inches(1.2))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)

    sb = slide.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.2), Inches(0.85))
    tf = sb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(20)
    r.font.color.rgb = RGBColor(203, 213, 225)

    mb = slide.shapes.add_textbox(Inches(0.9), Inches(5.85), Inches(11), Inches(0.4))
    tf = mb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = meta_line
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(148, 163, 184)

    # VOC summary slide
    if isinstance(voc_my, dict) or isinstance(voc_comp, dict):
        my_name = metadata.get("my_company") or voc_my.get("company", "자사")
        comp_name = metadata.get("competitor") or voc_comp.get("company", "경쟁사")
        bullets = []
        for label, data in [(my_name, voc_my), (comp_name, voc_comp)]:
            if isinstance(data, dict) and data:
                bullets.append(f"{label}: VOC {data.get('voc_count', 'N/A')}건 · 대표 감성 {data.get('overall_sentiment', 'N/A')} · 만족도 {data.get('satisfaction_score', 'N/A')}")
                if data.get("top_negatives"):
                    bullets.append(f"{label} 주요 불만: {', '.join(map(str, data.get('top_negatives', [])[:3]))}")
                if data.get("top_positives"):
                    bullets.append(f"{label} 긍정 포인트: {', '.join(map(str, data.get('top_positives', [])[:3]))}")
        if bullets:
            add_pptx_bullet_slide(prs, "VOC 분석 요약", compact_bullets(bullets, 7, 140))

    # Spec matrix table slide
    matrix = spec_comp.get("spec_matrix") if isinstance(spec_comp, dict) else None
    if isinstance(matrix, list) and matrix:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "제품 사양 비교 매트릭스"
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.name = "Malgun Gothic"
        r.font.color.rgb = RGBColor(30, 64, 175)

        rows = min(len(matrix), 8) + 1
        cols = 5
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(1.25), Inches(12.2), Inches(5.65))
        table = table_shape.table
        headers = ["항목", "자사", "경쟁사", "우위", "비고"]
        for c, h in enumerate(headers):
            table.cell(0, c).text = h
        for r_idx, item in enumerate(matrix[:8], start=1):
            values = [
                item.get("category", ""), item.get("my_value", ""), item.get("comp_value", ""),
                item.get("winner", ""), item.get("note", "")
            ]
            for c, val in enumerate(values):
                txt = strip_md_inline(str(val))
                table.cell(r_idx, c).text = txt[:75] + ("…" if len(txt) > 75 else "")
        for row in table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(9)
                    para.font.name = "Malgun Gothic"

    # Main report sections
    for section_title, lines in split_report_sections(report_md)[:9]:
        bullets = compact_bullets(lines, max_items=7)
        if bullets:
            add_pptx_bullet_slide(prs, section_title, bullets)

    # Recommendation slide
    rec = spec_comp.get("strategic_recommendation") if isinstance(spec_comp, dict) else ""
    if rec:
        add_pptx_bullet_slide(prs, "전략적 권고사항", compact_bullets([str(rec)], max_items=3, max_chars=160), accent_rgb=(124, 58, 237))

    # Closing slide
    add_pptx_bullet_slide(prs, "Appendix", [
        "본 덱은 앱에서 생성된 마크다운 보고서를 기반으로 자동 변환되었습니다.",
        "상세 표와 원문형 보고서는 DOCX/Markdown 파일을 함께 확인하세요.",
        f"생성 모델: {metadata.get('model', '')}",
    ], accent_rgb=(8, 145, 178))

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SIDEBAR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
with st.sidebar:
    st.markdown("## ⚙️ 시스템 설정")
    st.markdown("---")

    st.markdown("### 🔑 HuggingFace 인증")
    hf_token = st.text_input(
        "API Token (HF_TOKEN)", value=st.session_state.hf_token,
        type="password", placeholder="hf_xxxxxxxxxxxxxxxxxxxx",
        help="huggingface.co → Settings → Access Tokens 에서 발급",
    )
    if hf_token:
        st.session_state.hf_token = hf_token

    # ── 모델 선택 ──────────────────────────────────────────────────────────
    st.markdown("### 🤖 LLM 모델 선택")

    PRESET_MODELS = {
        "🌟 Google Gemma 4 26B (기본)":  "google/gemma-4-26B-A4B-it",
        "⚡ Zyphra ZAYA1 8B (경량·빠름)": "Zyphra/ZAYA1-8B",
        "✏️ 직접 입력":                   "__custom__",
    }

    # 현재 model_id 가 프리셋에 있으면 해당 항목, 없으면 직접 입력으로
    _reverse = {v: k for k, v in PRESET_MODELS.items() if v != "__custom__"}
    _default_label = _reverse.get(st.session_state.model_id, "✏️ 직접 입력")

    selected_label = st.selectbox(
        "모델 선택",
        options=list(PRESET_MODELS.keys()),
        index=list(PRESET_MODELS.keys()).index(_default_label),
        help="HuggingFace Inference API가 지원하는 모델만 동작합니다.",
    )

    selected_model_id = PRESET_MODELS[selected_label]

    if selected_model_id == "__custom__":
        custom_id = st.text_input(
            "모델 ID 직접 입력",
            value=(st.session_state.model_id
                   if st.session_state.model_id not in _reverse.values()
                   else ""),
            placeholder="예: meta-llama/Llama-3.1-8B-Instruct",
        )
        if custom_id.strip():
            st.session_state.model_id = custom_id.strip()
    else:
        st.session_state.model_id = selected_model_id

    # 현재 선택된 모델 표시
    _short = st.session_state.model_id.split("/")[-1]
    st.markdown(
        f'<div style="background:#0f172a;border:1px solid #334155;border-radius:8px;'
        f'padding:8px 12px;margin-top:4px;font-size:0.78rem;color:#94a3b8;">'
        f'적용 모델: <code style="color:#c4b5fd">{_short}</code></div>',
        unsafe_allow_html=True,
    )

    st.markdown("### 🌐 Inference Provider")
    PROVIDER_OPTIONS = {
        "자동 선택 (auto)": "auto",
        "Novita AI": "novita",
        "HF Inference": "hf-inference",
        "Together AI": "together",
        "DeepInfra": "deepinfra",
        "Nebius AI Studio": "nebius",
        "Fireworks AI": "fireworks-ai",
        "Groq": "groq",
    }
    provider_reverse = {v: k for k, v in PROVIDER_OPTIONS.items()}
    provider_label = st.selectbox(
        "Provider 선택",
        options=list(PROVIDER_OPTIONS.keys()),
        index=list(PROVIDER_OPTIONS.keys()).index(
            provider_reverse.get(st.session_state.get("hf_provider", "auto"), "자동 선택 (auto)")
        ),
        help="오류가 나면 auto → Novita AI 또는 HF Inference 등으로 바꿔 테스트하세요.",
    )
    st.session_state.hf_provider = PROVIDER_OPTIONS[provider_label]

    if st.session_state.hf_token:
        st.markdown('<span class="badge-ok">✅ API 토큰 설정됨</span>', unsafe_allow_html=True)
        if st.button("🧪 LLM 연결 테스트", use_container_width=True):
            with st.spinner("LLM 연결 테스트 중…"):
                test = call_llm("한국어로 한 문장만 답하세요: 연결 테스트 성공", max_tokens=80)
            if test.startswith("❌") or test.startswith("⚠️"):
                st.error(test)
            else:
                st.success("✅ 연결 성공")
                st.caption(test[:300])
    else:
        st.markdown('<span class="badge-warn">⚠️ 토큰 미설정</span>', unsafe_allow_html=True)

    st.markdown("---")
    st.markdown("### 📊 작업 현황")
    for label, done in [
        ("가이드 추출", st.session_state.guide_extracted),
        ("자사 VOC",    bool(st.session_state.voc_analyzed_my)),
        ("경쟁사 VOC",  bool(st.session_state.voc_analyzed_comp)),
        ("사양 비교",   bool(st.session_state.spec_comparison)),
        ("보고서",      bool(st.session_state.report_md)),
    ]:
        color = "#6ee7b7" if done else "#64748b"
        icon  = "✅" if done else "⬜"
        st.markdown(f'<div style="color:{color};font-size:0.9rem;padding:3px 0">{icon} {label}</div>', unsafe_allow_html=True)

    st.markdown("---")
    if st.button("🔄 전체 초기화", use_container_width=True):
        for k in list(DEFAULTS.keys()):
            if k not in ("hf_token","model_id","hf_provider"):
                st.session_state[k] = DEFAULTS[k]
        st.rerun()

    st.markdown('<div style="color:#475569;font-size:0.75rem;text-align:center">AI 벤치마킹 시스템 v2.0<br>Powered by HuggingFace</div>', unsafe_allow_html=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# HERO HEADER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
st.markdown("""
<div class="hero-header">
  <div class="hero-title">📊 AI 벤치마킹 보고서 자동 작성 시스템</div>
  <div class="hero-subtitle">벤치마킹 가이드 추출 → 실시간 VOC 수집 → 제품 사양 비교 → 전문 보고서 자동 생성</div>
  <div class="hero-badge">🤖 Powered by HuggingFace LLM</div>
</div>
""", unsafe_allow_html=True)

# 현재 선택된 모델을 히어로 아래 실시간 표시
st.markdown(
    f'<div style="text-align:center;margin:-12px 0 18px 0;font-size:0.8rem;color:#64748b;">' +
    f'적용 모델: <code style="color:#a78bfa;background:#1e1b4b;' +
    f'padding:2px 8px;border-radius:4px;">{st.session_state.model_id}</code></div>',
    unsafe_allow_html=True,
)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MAIN TABS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
tab1, tab2, tab3, tab4, tab5 = st.tabs([
    "📋 Step 1 · 가이드 추출",
    "🗣️ Step 2 · VOC 수집",
    "🔬 Step 3 · 사양 비교",
    "📝 Step 4 · 보고서 생성",
    "📖 도움말",
])

# ──────────────────────────────────────────────────────────────────────────────
# TAB 1 – GUIDE EXTRACTION
# ──────────────────────────────────────────────────────────────────────────────
with tab1:
    st.markdown("## 📋 벤치마킹 보고서 작성 기준 가이드 추출")
    st.markdown("기존 벤치마킹 보고서 파일을 업로드하면 AI가 공통 구조·핵심 항목·평가 기준을 자동으로 분석하여 표준 가이드를 생성합니다.")

    with st.expander("📁 보고서 파일 업로드 (PDF · DOCX · TXT · MD)", expanded=True):
        uploaded = st.file_uploader(
            "파일을 드래그하거나 클릭하여 선택 (다중 선택 가능)",
            accept_multiple_files=True,
            type=["pdf","docx","doc","txt","md"],
        )
        if uploaded:
            st.success(f"✅ {len(uploaded)}개 파일 업로드 완료")
            cols = st.columns(3)
            for i, f in enumerate(uploaded):
                cols[i%3].markdown(
                    f'<div class="step-card"><h4>📄 {f.name}</h4><p>{f.size:,} bytes</p></div>',
                    unsafe_allow_html=True,
                )

    c1, c2, c3 = st.columns([2,2,1])
    with c1:
        if st.button("🚀 AI로 가이드 자동 추출", type="primary", use_container_width=True, disabled=not uploaded):
            if not st.session_state.hf_token:
                st.error("❌ HuggingFace API 토큰을 먼저 입력해주세요.")
            else:
                prog = st.progress(0, "파일 텍스트 추출 중…")
                texts = []
                for i, f in enumerate(uploaded):
                    prog.progress((i+1)/len(uploaded), f"추출 중: {f.name}")
                    t = extract_file(f)
                    if t.strip():
                        texts.append(f"[파일: {f.name}]\n{t}")
                if texts:
                    prog.progress(1.0, "LLM으로 가이드 생성 중…")
                    with st.spinner("🤖 AI가 보고서 패턴을 분석하고 있습니다…"):
                        guide = build_guide_from_files(texts)
                    st.session_state.benchmark_guide = guide
                    st.session_state.guide_extracted = True
                    prog.empty()
                    st.success(f"✅ {len(texts)}개 파일 분석 완료! 가이드가 생성되었습니다.")
                else:
                    st.error("파일에서 텍스트를 추출할 수 없었습니다.")

    with c2:
        if st.button("📝 기본 표준 가이드 사용", use_container_width=True):
            st.session_state.benchmark_guide = DEFAULT_GUIDE
            st.session_state.guide_extracted = True
            st.success("✅ 표준 가이드가 적용되었습니다.")

    with c3:
        if st.session_state.benchmark_guide:
            st.download_button(
                "💾 JSON 저장",
                json.dumps(st.session_state.benchmark_guide, ensure_ascii=False, indent=2),
                "benchmark_guide.json", "application/json", use_container_width=True,
            )

    if st.session_state.benchmark_guide:
        st.markdown("---")
        g = st.session_state.benchmark_guide
        st.markdown(f"### 📌 {g.get('title','가이드')}")
        if g.get("summary"):
            st.info(g["summary"])
        if g.get("llm_notes"):
            st.text_area("LLM 분석 원문", g["llm_notes"], height=200)

        sections = g.get("main_sections", [])
        if sections:
            st.markdown("#### 📂 주요 섹션 구성")
            cols_g = st.columns(2)
            for idx, sec in enumerate(sections):
                with cols_g[idx%2]:
                    with st.expander(f"{sec.get('id','')+'.' if 'id' in sec else ''} {sec.get('name','섹션')}"):
                        if sec.get("items"):
                            st.markdown("**핵심 항목:**")
                            for it in sec["items"]:
                                st.markdown(f"  - {it}")
                        if sec.get("criteria"):
                            st.markdown("**평가 기준:**")
                            for cr in sec["criteria"]:
                                st.markdown(f"  - {cr}")

        r1, r2 = st.columns(2)
        with r1:
            if g.get("scoring_dimensions"):
                st.markdown("#### 🎯 평가 차원")
                for d in g["scoring_dimensions"]: st.markdown(f"  - {d}")
        with r2:
            if g.get("analysis_methods"):
                st.markdown("#### 🔬 분석 방법론")
                for m in g["analysis_methods"]: st.markdown(f"  - {m}")


# ──────────────────────────────────────────────────────────────────────────────
# TAB 2 – VOC COLLECTION
# ──────────────────────────────────────────────────────────────────────────────
with tab2:
    st.markdown("## 🗣️ 실시간 VOC 수집 및 AI 감성 분석")
    st.markdown("웹 검색·뉴스·RSS를 함께 수집하고, 중복 제거·카테고리 분류·감성 힌트·LLM 심층 분석까지 자동 수행합니다.")

    v1, v2 = st.columns(2)
    with v1:
        st.markdown("### 🏢 자사")
        my_co_v = st.text_input("회사명", key="v_my_co", placeholder="예: 삼성전자")
        my_pr_v = st.text_input("제품/서비스명", key="v_my_pr", placeholder="예: Galaxy S25 Ultra")
    with v2:
        st.markdown("### 🏭 경쟁사")
        cp_co_v = st.text_input("회사명", key="v_cp_co", placeholder="예: Apple")
        cp_pr_v = st.text_input("제품/서비스명", key="v_cp_pr", placeholder="예: iPhone 16 Pro Max")

    st.markdown("### ⚙️ 수집 옵션")
    opt1, opt2, opt3, opt4 = st.columns(4)
    with opt1:
        n_voc = st.slider("쿼리당 수집 건수", 4, 25, 10)
    with opt2:
        depth_label = st.selectbox("수집 깊이", ["기본", "표준", "심층"], index=1)
        depth_map = {"기본": "basic", "표준": "standard", "심층": "deep"}
    with opt3:
        time_label = st.selectbox("검색 기간", ["전체", "최근 1일", "최근 1주", "최근 1개월", "최근 1년"], index=3)
        time_map = {"전체": None, "최근 1일": "d", "최근 1주": "w", "최근 1개월": "m", "최근 1년": "y"}
    with opt4:
        region = st.selectbox("검색 지역", ["kr-kr", "us-en", "wt-wt"], index=0)

    source_cols = st.columns(4)
    with source_cols[0]:
        src_web = st.checkbox("일반 웹", value=True)
    with source_cols[1]:
        src_news = st.checkbox("뉴스 검색", value=True)
    with source_cols[2]:
        src_rss = st.checkbox("Google News RSS", value=True)
    with source_cols[3]:
        fetch_pages = st.checkbox("본문 일부 보강", value=False, help="검색 snippet이 짧을 때 URL 본문 일부를 가져옵니다. 느려질 수 있습니다.")

    selected_sources = []
    if src_web: selected_sources.append("web")
    if src_news: selected_sources.append("news")
    if src_rss: selected_sources.append("rss")

    if st.button("🔍 VOC 수집 및 분석 시작", type="primary", use_container_width=True, disabled=not (my_co_v and cp_co_v and selected_sources)):
        if not st.session_state.hf_token:
            st.error("❌ HuggingFace API 토큰을 먼저 입력해주세요.")
        else:
            st.session_state.update(my_company=my_co_v, my_product=my_pr_v, competitor=cp_co_v, comp_product=cp_pr_v)
            p1, p2 = st.columns(2)
            with p1:
                with st.spinner(f"🔍 {my_co_v} VOC 멀티소스 수집 중…"):
                    raw_my = collect_voc(
                        my_co_v, my_pr_v, n_voc,
                        sources=selected_sources,
                        depth=depth_map[depth_label],
                        fetch_pages=fetch_pages,
                        region=region,
                        timelimit=time_map[time_label],
                    )
                    st.session_state.voc_my = raw_my
                with st.spinner(f"🤖 {my_co_v} VOC 심층 분석 중…"):
                    st.session_state.voc_analyzed_my = analyze_voc(my_co_v, my_pr_v, raw_my)
                st.success(f"✅ {my_co_v}: {len(raw_my)}건 수집 완료")
            with p2:
                with st.spinner(f"🔍 {cp_co_v} VOC 멀티소스 수집 중…"):
                    raw_cp = collect_voc(
                        cp_co_v, cp_pr_v, n_voc,
                        sources=selected_sources,
                        depth=depth_map[depth_label],
                        fetch_pages=fetch_pages,
                        region=region,
                        timelimit=time_map[time_label],
                    )
                    st.session_state.voc_comp = raw_cp
                with st.spinner(f"🤖 {cp_co_v} VOC 심층 분석 중…"):
                    st.session_state.voc_analyzed_comp = analyze_voc(cp_co_v, cp_pr_v, raw_cp)
                st.success(f"✅ {cp_co_v}: {len(raw_cp)}건 수집 완료")

    def render_distribution(title, dist: dict):
        if not isinstance(dist, dict) or not dist:
            return
        df = pd.DataFrame([{"항목": str(k), "건수": v} for k, v in dist.items()])
        df["건수"] = pd.to_numeric(df["건수"], errors="coerce").fillna(0).astype(int)
        df = df[df["건수"] > 0]
        if df.empty:
            return
        st.markdown(f"**{title}**")
        st.bar_chart(df.set_index("항목"))

    def render_raw_voc_tools(raw, label):
        if not raw:
            return
        df = pd.DataFrame(raw)
        view_cols = [c for c in ["title", "source_domain", "channel", "source_type", "category", "sentiment_hint", "relevance_score", "date", "href"] if c in df.columns]
        with st.expander(f"📄 정제 VOC 테이블 · {label} ({len(df)}건)"):
            st.dataframe(df[view_cols], use_container_width=True, height=320)
            csv = df.to_csv(index=False).encode("utf-8-sig")
            st.download_button(
                "⬇️ CSV 다운로드",
                data=csv,
                file_name=f"voc_{re.sub(r'[^0-9A-Za-z가-힣_-]+', '_', label)}.csv",
                mime="text/csv",
                use_container_width=True,
            )

    def render_voc_panel(col, data, raw, label):
        with col:
            score     = data.get("satisfaction_score","N/A") if isinstance(data, dict) else "N/A"
            sentiment = data.get("overall_sentiment","N/A") if isinstance(data, dict) else "N/A"
            count     = data.get("voc_count", len(raw)) if isinstance(data, dict) else len(raw)
            confidence = data.get("confidence", "N/A") if isinstance(data, dict) else "N/A"
            st.markdown(f"#### 🏢 {label}")
            m1,m2,m3,m4 = st.columns(4)
            m1.metric("만족도",   f"{score}/10" if score!="N/A" else "N/A")
            m2.metric("대표 감성", sentiment)
            m3.metric("수집 VOC", f"{count}건")
            m4.metric("신뢰도", confidence)

            if isinstance(data, dict):
                c1, c2 = st.columns(2)
                with c1:
                    render_distribution("카테고리 분포", data.get("category_distribution", {}))
                with c2:
                    render_distribution("감성 분포", data.get("sentiment_distribution", {}))

                if data.get("raw"):
                    st.warning("LLM JSON 파싱 실패 또는 API 오류로 로컬 통계 fallback을 표시합니다.")
                    st.text_area("LLM 원문", data["raw"][:1500], height=140)

                for title, key, icon in [
                    ("주요 긍정 포인트","top_positives","✅"),
                    ("주요 불만 사항","top_negatives","❌"),
                    ("개선 요청","improvement_requests","💡"),
                    ("트렌딩 이슈","trending_topics","🔥"),
                    ("칭찬받는 기능","key_features_praised","⭐"),
                    ("비판받는 기능","key_features_criticized","⚠️"),
                ]:
                    if data.get(key):
                        with st.expander(f"{icon} {title}"):
                            for item in data[key]: st.markdown(f"- {item}")

                if data.get("issue_matrix"):
                    with st.expander("🧩 이슈 매트릭스"):
                        st.dataframe(pd.DataFrame(data["issue_matrix"]), use_container_width=True)

            with st.expander(f"🔎 원시 VOC 미리보기 ({len(raw)}건 중 5건)"):
                for i, item in enumerate(raw[:5]):
                    st.markdown(
                        f'<div class="voc-card"><div class="voc-title">{i+1}. {item.get("title","N/A")}</div>'
                        f'<div class="voc-body">[{item.get("channel", "웹")}/{item.get("category", "기타")}/{item.get("sentiment_hint", "중립")}] '
                        f'{(item.get("text") or item.get("body") or item.get("snippet", ""))[:360]}</div></div>',
                        unsafe_allow_html=True,
                    )
            render_raw_voc_tools(raw, label)

    if st.session_state.voc_analyzed_my or st.session_state.voc_analyzed_comp:
        st.markdown("---"); st.markdown("### 📊 VOC 분석 결과")
        d1, d2 = st.columns(2)
        render_voc_panel(d1, st.session_state.voc_analyzed_my,  st.session_state.voc_my,  f"{st.session_state.my_company} {st.session_state.my_product}")
        render_voc_panel(d2, st.session_state.voc_analyzed_comp,st.session_state.voc_comp, f"{st.session_state.competitor} {st.session_state.comp_product}")


# ──────────────────────────────────────────────────────────────────────────────
# TAB 3 – SPEC COMPARISON
# ──────────────────────────────────────────────────────────────────────────────
with tab3:
    st.markdown("## 🔬 제품 사양 및 성능 비교 분석")
    s1, s2 = st.columns(2)
    with s1:
        st.markdown("### 🏢 자사 제품")
        sp_my_co = st.text_input("회사명", value=st.session_state.my_company,  key="sp_my_co")
        sp_my_pr = st.text_input("제품명", value=st.session_state.my_product,  key="sp_my_pr")
        my_spec_txt = st.text_area("제품 사양/특징 (직접 입력)", value=st.session_state.spec_my_raw, height=220,
            placeholder="예:\n- AP: Snapdragon 8 Elite\n- 배터리: 5000mAh\n- 카메라: 200MP", key="sp_my_spec")
        st.session_state.spec_my_raw = my_spec_txt
    with s2:
        st.markdown("### 🏭 경쟁사 제품")
        sp_cp_co = st.text_input("회사명",   value=st.session_state.competitor,   key="sp_cp_co")
        sp_cp_pr = st.text_input("제품명",   value=st.session_state.comp_product, key="sp_cp_pr")
        cp_spec_txt = st.text_area("제품 사양/특징 (직접 입력)", value=st.session_state.spec_comp_raw, height=220,
            placeholder="예:\n- AP: A18 Pro\n- 배터리: 4422mAh\n- 카메라: 48MP 트리플", key="sp_cp_spec")
        st.session_state.spec_comp_raw = cp_spec_txt

    b1, b2 = st.columns(2)
    with b1:
        if st.button("🌐 웹에서 사양 자동 수집", use_container_width=True, disabled=not (sp_my_co and sp_cp_co)):
            with st.spinner("웹 검색으로 사양 수집 중…"):
                my_res  = collect_specs(sp_my_co, sp_my_pr)
                cp_res  = collect_specs(sp_cp_co, sp_cp_pr)
                st.session_state.spec_my_raw   = "\n".join(f"[{r.get('title','')}] {r.get('body',r.get('snippet',''))}" for r in my_res[:6])[:3000]
                st.session_state.spec_comp_raw = "\n".join(f"[{r.get('title','')}] {r.get('body',r.get('snippet',''))}" for r in cp_res[:6])[:3000]
                st.success("✅ 웹 사양 수집 완료!"); st.rerun()
    with b2:
        if st.button("⚡ AI 사양 비교 분석 실행", type="primary", use_container_width=True):
            if not st.session_state.hf_token:
                st.error("❌ HuggingFace API 토큰을 먼저 입력해주세요.")
            elif not (sp_my_co and sp_cp_co):
                st.error("회사명을 모두 입력해주세요.")
            else:
                st.session_state.update(my_company=sp_my_co, my_product=sp_my_pr, competitor=sp_cp_co, comp_product=sp_cp_pr)
                with st.spinner("🤖 AI가 사양을 비교 분석 중…"):
                    result = compare_specs(sp_my_co, sp_my_pr, sp_cp_co, sp_cp_pr,
                                           st.session_state.spec_my_raw, st.session_state.spec_comp_raw)
                    st.session_state.spec_comparison = result
                st.success("✅ 사양 비교 완료!")

    if st.session_state.spec_comparison:
        st.markdown("---")
        comp = st.session_state.spec_comparison
        if comp.get("raw"):
            st.text_area("LLM 분석 원문", comp["raw"][:2000], height=200)
        else:
            st.markdown("### 📊 비교 분석 결과")
            if comp.get("overall_summary"): st.info(f"**종합 평가:** {comp['overall_summary']}")
            if comp.get("overall_winner"):  st.metric("🏆 종합 우위", comp["overall_winner"])

            if comp.get("spec_matrix"):
                st.markdown("#### 📋 상세 사양 비교 매트릭스")
                rows = []
                for row in comp["spec_matrix"]:
                    w = row.get("winner","")
                    rows.append({
                        "카테고리": row.get("category",""),
                        "중요도":   row.get("importance",""),
                        f"{sp_my_co}": row.get("my_value",""),
                        f"{sp_cp_co}": row.get("comp_value",""),
                        "우위": f"✅ {w}" if sp_my_co in w else (f"❌ {w}" if w!="동등" else "🟡 동등"),
                        "비고": row.get("note",""),
                    })
                st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)

            sg1, sg2 = st.columns(2)
            with sg1:
                for title, key in [("강점","my_strengths"),("고유 기능","my_unique_only"),("취약 영역","my_gaps")]:
                    if comp.get(key):
                        st.markdown(f"#### {sp_my_co} {title}")
                        for item in comp[key]: st.markdown(f"- {item}")
            with sg2:
                for title, key in [("강점","comp_strengths"),("고유 기능","comp_unique_only"),("취약 영역","comp_gaps")]:
                    if comp.get(key):
                        st.markdown(f"#### {sp_cp_co} {title}")
                        for item in comp[key]: st.markdown(f"- {item}")

            if comp.get("differentiation_score"):
                st.markdown("#### 🎯 차별화 점수 비교")
                sdf = pd.DataFrame([
                    {"평가 차원": d, f"{sp_my_co} 점수": v.get("my",0), f"{sp_cp_co} 점수": v.get("comp",0)}
                    for d, v in comp["differentiation_score"].items()
                ])
                st.dataframe(sdf, use_container_width=True, hide_index=True)

            if comp.get("strategic_recommendation"):
                st.success(f"💡 **전략적 권고:** {comp['strategic_recommendation']}")


# ──────────────────────────────────────────────────────────────────────────────
# TAB 4 – REPORT GENERATION
# ──────────────────────────────────────────────────────────────────────────────
with tab4:
    st.markdown("## 📝 벤치마킹 보고서 자동 생성")
    r1c,r2c,r3c,r4c = st.columns(4)
    r1c.metric("가이드",   "✅" if st.session_state.guide_extracted       else "⚠️")
    r2c.metric("자사 VOC", "✅" if st.session_state.voc_analyzed_my      else "⚠️")
    r3c.metric("경쟁사 VOC","✅" if st.session_state.voc_analyzed_comp   else "⚠️")
    r4c.metric("사양 비교","✅" if st.session_state.spec_comparison       else "⚠️")

    with st.expander("⚙️ 보고서 기본 정보"):
        re1, re2 = st.columns(2)
        with re1:
            r_my_co = st.text_input("자사명",   value=st.session_state.my_company,  key="r_my_co")
            r_my_pr = st.text_input("자사 제품", value=st.session_state.my_product,  key="r_my_pr")
        with re2:
            r_cp_co = st.text_input("경쟁사명",   value=st.session_state.competitor,   key="r_cp_co")
            r_cp_pr = st.text_input("경쟁사 제품", value=st.session_state.comp_product, key="r_cp_pr")

    if st.button("🚀 벤치마킹 보고서 생성 (LLM 기반)", type="primary", use_container_width=True, disabled=not st.session_state.hf_token):
        if not (r_my_co and r_cp_co):
            st.error("자사명과 경쟁사명을 입력해주세요.")
        else:
            if not st.session_state.benchmark_guide:
                st.session_state.benchmark_guide = DEFAULT_GUIDE
                st.session_state.guide_extracted = True
            with st.spinner("🤖 AI가 벤치마킹 보고서를 작성 중입니다… (2~5분 소요)"):
                report = generate_report(
                    st.session_state.benchmark_guide,
                    r_my_co, r_my_pr, r_cp_co, r_cp_pr,
                    st.session_state.voc_analyzed_my, st.session_state.voc_analyzed_comp,
                    st.session_state.spec_comparison,
                )
                st.session_state.report_md = report
            st.success("✅ 보고서 생성 완료!")

    if st.session_state.report_md:
        st.markdown("---")
        st.markdown("### 📄 생성된 보고서")
        st.markdown(st.session_state.report_md)
        st.markdown("---")

        today_str = datetime.now().strftime("%Y%m%d_%H%M")
        raw_fname = f"benchmarking_{st.session_state.my_company}_vs_{st.session_state.competitor}_{today_str}"
        fname = safe_filename(raw_fname)
        export_meta = {
            "title": f"벤치마킹 보고서: {st.session_state.my_company or r_my_co} vs {st.session_state.competitor or r_cp_co}",
            "subtitle": f"대상 제품: {st.session_state.my_product or r_my_pr} vs {st.session_state.comp_product or r_cp_pr}",
            "meta_line": f"작성일: {datetime.now().strftime('%Y-%m-%d %H:%M')} | 모델: {st.session_state.model_id}",
            "model": st.session_state.model_id,
            "my_company": st.session_state.my_company or r_my_co,
            "competitor": st.session_state.competitor or r_cp_co,
        }
        dl1, dl2, dl3, dl4, dl5 = st.columns(5)

        with dl1:
            st.download_button("📥 Markdown", st.session_state.report_md,
                               f"{fname}.md", "text/markdown", use_container_width=True)
        with dl2:
            html_doc = f"""<!DOCTYPE html><html lang="ko"><head><meta charset="UTF-8">
<title>벤치마킹 보고서</title>
<style>body{{font-family:'Malgun Gothic',sans-serif;max-width:1100px;margin:40px auto;padding:0 24px;color:#1e293b;line-height:1.8}}
h1{{color:#1e40af;border-bottom:3px solid #1e40af;padding-bottom:12px}}h2{{color:#2a5298}}
table{{border-collapse:collapse;width:100%;margin:16px 0}}th{{background:#1e40af;color:white;padding:10px 14px}}
td{{padding:9px 14px;border-bottom:1px solid #e2e8f0}}tr:nth-child(even){{background:#f8fafc}}</style>
</head><body>{st.session_state.report_md.replace(chr(10),'<br>')}</body></html>"""
            st.download_button("📥 HTML", html_doc,
                               f"{fname}.html", "text/html", use_container_width=True)
        with dl3:
            try:
                docx_bytes = build_docx_report_bytes(st.session_state.report_md, export_meta)
                st.download_button("📥 Word DOCX", docx_bytes,
                                   f"{fname}.docx",
                                   "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                   use_container_width=True)
            except Exception as e:
                st.error(f"DOCX 생성 오류: {e}")
        with dl4:
            try:
                pptx_bytes = build_pptx_report_bytes(
                    st.session_state.report_md,
                    export_meta,
                    st.session_state.voc_analyzed_my,
                    st.session_state.voc_analyzed_comp,
                    st.session_state.spec_comparison,
                )
                st.download_button("📥 PowerPoint PPTX", pptx_bytes,
                                   f"{fname}.pptx",
                                   "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                   use_container_width=True)
            except Exception as e:
                st.error(f"PPTX 생성 오류: {e}")
        with dl5:
            full_json = json.dumps({
                "metadata": {"generated_at": datetime.now().isoformat(), "model": st.session_state.model_id,
                             "my_company": st.session_state.my_company, "competitor": st.session_state.competitor},
                "guide": st.session_state.benchmark_guide,
                "voc_my": st.session_state.voc_analyzed_my,
                "voc_comp": st.session_state.voc_analyzed_comp,
                "spec_comparison": st.session_state.spec_comparison,
                "report_md": st.session_state.report_md,
            }, ensure_ascii=False, indent=2)
            st.download_button("📥 전체 JSON", full_json,
                               f"{fname}_full.json", "application/json", use_container_width=True)


# ──────────────────────────────────────────────────────────────────────────────
# TAB 5 – HELP
# ──────────────────────────────────────────────────────────────────────────────
with tab5:
    st.markdown("## 📖 사용 가이드 및 시스템 정보")
    with st.expander("🚀 빠른 시작 가이드", expanded=True):
        st.markdown("""
### 사전 준비
1. **HuggingFace 계정 생성** → [huggingface.co](https://huggingface.co)
2. **API Token 발급** → Settings → Access Tokens → New Token (Read 권한)
3. **모델 접근 동의** → `google/gemma-4-26B-A4B-it` 모델 페이지에서 라이선스 동의
4. 사이드바에 토큰 입력 후 사용 시작

### 4단계 워크플로우

| 단계 | 작업 | 소요 시간 |
|------|------|-----------|
| Step 1 | 보고서 파일 업로드 → AI 가이드 추출 | 1~3분 |
| Step 2 | 자사·경쟁사 입력 → 멀티소스 VOC 수집·분류·분석 | 2~7분 |
| Step 3 | 사양 입력/수집 → AI 비교 분석 | 1~2분 |
| Step 4 | 버튼 클릭 → 전문 보고서 자동 생성 | 2~5분 |
""")
    with st.expander("⚙️ 기술 스택"):
        st.markdown("""
| 구성 요소 | 기술 |
|-----------|------|
| UI | Streamlit |
| LLM | google/gemma-4-26B-A4B-it (HuggingFace Inference API) |
| 웹/VOC 검색 | DuckDuckGo Web · DuckDuckGo News · Google News RSS |
| 파일 파싱 | PyPDF2 · python-docx |
| 웹 스크래핑 | requests · BeautifulSoup4 · lxml |
| 데이터 처리 | pandas |
""")
    with st.expander("❓ 트러블슈팅"):
        st.markdown("""
| 증상 | 해결 방법 |
|------|-----------|
| 401 Unauthorized | HF 토큰 재확인·재발급 |
| Model not found | HF 모델 페이지에서 라이선스 동의 |
| not supported for task text-generation | Provider/model은 conversational만 지원하므로 chat-completion 방식으로 호출해야 함 |
| Provider 오류 | 사이드바 Provider를 auto, Novita AI, HF Inference 등으로 바꿔 연결 테스트 |
| VOC 수집 0건 | 소스를 웹+뉴스+RSS로 모두 켜고, 검색 기간을 전체 또는 최근 1년으로 변경 |
| PDF 텍스트 없음 | OCR 전처리 후 재업로드 |
""")
    st.markdown(f"""
---
<div style="text-align:center;color:#475569;font-size:0.85rem">
AI 벤치마킹 보고서 자동 작성 시스템 v2.0<br>
HuggingFace LLM · Multi-source VOC Collector · Streamlit<br>
현재 모델: <code>{st.session_state.model_id}</code>
</div>
""", unsafe_allow_html=True)
